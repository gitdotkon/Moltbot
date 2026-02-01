#!/usr/bin/env python3
"""
PPT Generator with Search Integration
从GitHub、StackOverflow等开源网站搜索并生成PPT
"""

import os
import json
import requests
from datetime import datetime
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 配置
PPT_DIR = os.environ.get('PPT_DIR', os.path.expanduser('~/.clawnn/ppt'))
TEMPLATE_DIR = os.path.join(PPT_DIR, 'templates')
os.makedirs(PPT_DIR, exist_ok=True)
os.makedirs(TEMPLATE_DIR, exist_ok=True)


def search_github(keyword, limit=10):
    """搜索GitHub热门项目"""
    try:
        url = f"https://api.github.com/search/repositories?q={keyword}&sort=stars&per_page={limit}"
        resp = requests.get(url, timeout=10)
        if resp.status_code == 200:
            data = resp.json()
            return [{
                'source': 'GitHub',
                'title': item.get('full_name', ''),
                'description': item.get('description', '')[:200],
                'url': item.get('html_url', ''),
                'stars': item.get('stargazers_count', 0),
                'language': item.get('language', '')
            } for item in data.get('items', [])]
    except Exception as e:
        print(f"GitHub搜索失败: {e}")
    return []


def search_stackoverflow(keyword, limit=10):
    """搜索StackOverflow问答"""
    try:
        url = f"https://stackoverflow.com/search?q={keyword}"
        headers = {'User-Agent': 'Mozilla/5.0'}
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code == 200:
            soup = BeautifulSoup(resp.text, 'html.parser')
            results = []
            for item in soup.select('.s-post-summary')[:limit]:
                title_elem = item.select_one('.s-post-summary--content--title a')
                votes_elem = item.select_one('.s-post-summary--stats .vote-count-post')
                if title_elem:
                    results.append({
                        'source': 'StackOverflow',
                        'title': title_elem.get_text(strip=True),
                        'description': '',
                        'url': 'https://stackoverflow.com' + title_elem.get('href', ''),
                        'votes': votes_elem.get_text(strip=True) if votes_elem else '0'
                    })
            return results
    except Exception as e:
        print(f"StackOverflow搜索失败: {e}")
    return []


def search_all(keyword, sources=['github', 'stackoverflow'], limit=10):
    """综合搜索"""
    all_results = []
    
    if 'github' in sources:
        all_results.extend(search_github(keyword, limit))
    
    if 'stackoverflow' in sources:
        all_results.extend(search_stackoverflow(keyword, limit))
    
    # 按stars/votes排序
    all_results.sort(key=lambda x: x.get('stars', 0) or int(x.get('votes', 0)), reverse=True)
    
    return all_results[:limit]


def create_ppt_from_search(keyword, results, output_file=None):
    """从搜索结果生成PPT"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"📚 {keyword} 资料汇总"
    subtitle.text = f"共找到 {len(results)} 条资源\n{datetime.now().strftime('%Y-%m-%d')}"
    
    # 内容页
    content_layout = prs.slide_layouts[1]
    
    # GitHub 项目页
    github_results = [r for r in results if r['source'] == 'GitHub']
    if github_results:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = "🐙 GitHub 热门项目"
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for i, r in enumerate(github_results[:5], 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {r['title']}"
            p.font.bold = True
            p.font.size = Pt(14)
            
            if r.get('description'):
                p2 = tf.add_paragraph()
                p2.text = f"   ⭐ {r['stars']} | {r['description'][:80]}..."
                p2.font.size = Pt(12)
            
            p3 = tf.add_paragraph()
            p3.text = f"   🔗 {r['url']}"
            p3.font.size = Pt(11)
            p3.font.color.rgb = 0x0000FF
            
            tf.add_paragraph()  # 空行
    
    # StackOverflow 问答页
    so_results = [r for r in results if r['source'] == 'StackOverflow']
    if so_results:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = "💬 StackOverflow 热门问答"
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for i, r in enumerate(so_results[:5], 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {r['title']}"
            p.font.bold = True
            p.font.size = Pt(14)
            
            p2 = tf.add_paragraph()
            p2.text = f"   👍 {r.get('votes', 0)} 票"
            p2.font.size = Pt(12)
            
            p3 = tf.add_paragraph()
            p3.text = f"   🔗 {r['url']}"
            p3.font.size = Pt(11)
            p3.font.color.rgb = 0x0000FF
            
            tf.add_paragraph()
    
    # 保存
    if not output_file:
        output_file = os.path.join(PPT_DIR, f'{keyword.replace(" ", "_")}_resources.pptx')
    
    prs.save(output_file)
    return output_file


def create_basic_ppt(title, outline_lines=None, slides_count=5, output_file=None):
    """创建基础PPT"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    
    # 内容页
    content_layout = prs.slide_layouts[1]
    
    lines = outline_lines or [f"第{i+1}页内容" for i in range(slides_count)]
    
    for i, line in enumerate(lines[:slides_count], 1):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = f"第 {i} 部分"
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = line
        
        # 添加默认要点
        tf.add_paragraph().text = "• 关键要点一"
        tf.add_paragraph().text = "• 关键要点二"
        tf.add_paragraph().text = "• 关键要点三"
    
    # 保存
    if not output_file:
        output_file = os.path.join(PPT_DIR, f'{title.replace(" ", "_")}.pptx')
    
    prs.save(output_file)
    return output_file


def update_templates():
    """更新模板列表"""
    templates = {
        "templates": [
            {"name": "modern", "style": "简约现代", "slides": 10},
            {"name": "business", "style": "商务正式", "slides": 15},
            {"name": "tech", "style": "技术风格", "slides": 12},
            {"name": "creative", "style": "创意风格", "slides": 8}
        ],
        "last_update": datetime.now().isoformat(),
        "sources": ["GitHub", "StackOverflow", "Dev.to", "Medium"]
    }
    
    template_file = os.path.join(TEMPLATE_DIR, 'templates.json')
    with open(template_file, 'w', encoding='utf-8') as f:
        json.dump(templates, f, ensure_ascii=False, indent=2)
    
    print(f"✅ 模板已更新: {template_file}")
    return template_file


if __name__ == '__main__':
    import sys
    
    if len(sys.argv) < 2:
        print("用法:")
        print("  python3 ppt_gen.py create <标题> [--outline '第一\\n第二\\n第三']")
        print("  python3 ppt_gen.py search <关键词> [--sources github,stackoverflow] [--limit 10]")
        print("  python3 ppt_gen.py update-templates")
        sys.exit(1)
    
    command = sys.argv[1]
    
    if command == 'create':
        title = sys.argv[2] if len(sys.argv) > 2 else "演示文稿"
        outline = ''
        output = None
        slides = 5
        
        i = 3
        while i < len(sys.argv):
            if sys.argv[i] == '--outline':
                outline = sys.argv[i+1].replace('\\n', '\n')
                i += 2
            elif sys.argv[i] == '--output':
                output = sys.argv[i+1]
                i += 2
            elif sys.argv[i] == '--slides':
                slides = int(sys.argv[i+1])
                i += 2
            else:
                i += 1
        
        outline_lines = outline.split('\n') if outline else None
        result = create_basic_ppt(title, outline_lines, slides, output)
        print(f"✅ PPT已生成: {result}")
    
    elif command == 'search':
        keyword = sys.argv[2] if len(sys.argv) > 2 else ""
        sources = ['github', 'stackoverflow']
        limit = 10
        output = None
        
        i = 3
        while i < len(sys.argv):
            if sys.argv[i] == '--sources':
                sources = sys.argv[i+1].split(',')
                i += 2
            elif sys.argv[i] == '--limit':
                limit = int(sys.argv[i+1])
                i += 2
            elif sys.argv[i] == '--output':
                output = sys.argv[i+1]
                i += 2
            else:
                i += 1
        
        if not keyword:
            print("请输入搜索关键词")
            sys.exit(1)
        
        print(f"🔍 搜索: {keyword}")
        results = search_all(keyword, sources, limit)
        
        if results:
            print(f"\n找到 {len(results)} 条结果:")
            for i, r in enumerate(results[:5], 1):
                print(f"  {i}. [{r['source']}] {r['title'][:50]}")
            
            result = create_ppt_from_search(keyword, results, output)
            print(f"\n✅ PPT已生成: {result}")
        else:
            print("未找到结果")
    
    elif command == 'update-templates':
        update_templates()
    
    else:
        print(f"未知命令: {command}")
